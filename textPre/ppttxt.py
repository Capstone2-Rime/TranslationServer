from pptx import Presentation
import re

def read_ppt(ppt_file_path):
    text_runs = []
    prs = Presentation(ppt_file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    text_runs = ' '.join(text_runs).split()
    text = ' '.join(text_runs)
    return text